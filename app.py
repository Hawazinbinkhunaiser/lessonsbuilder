import streamlit as st
import anthropic
import requests
import json
import os
from typing import List, Dict, Optional
import time
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_THEME_COLOR
import io
import base64
from PIL import Image, ImageDraw, ImageFont, ImageFilter, ImageEnhance
import tempfile
import subprocess
import numpy as np
import zipfile
import re
from datetime import datetime

# Try to import MoviePy with fallback
try:
    from moviepy.editor import ImageSequenceClip, AudioFileClip, CompositeVideoClip, concatenate_videoclips, TextClip
    MOVIEPY_AVAILABLE = True
except ImportError as e:
    MOVIEPY_AVAILABLE = False

# Configure page
st.set_page_config(
    page_title="Enhanced AI Lesson Generator",
    page_icon="🚀",
    layout="wide",
    initial_sidebar_state="expanded"
)

# Enhanced Custom CSS with sophisticated design
st.markdown("""
<style>
    @import url('https://fonts.googleapis.com/css2?family=Inter:wght@300;400;500;600;700&family=Playfair+Display:wght@400;700&family=Space+Grotesk:wght@300;400;500;600&display=swap');
    
    /* Global Styles */
    .main > div {
        padding-top: 2rem;
    }
    
    /* Enhanced color palette */
    :root {
        --primary-purple: #6366f1;
        --primary-purple-dark: #4f46e5;
        --secondary-teal: #14b8a6;
        --accent-coral: #f97316;
        --accent-rose: #ec4899;
        --neutral-dark: #1f2937;
        --neutral-medium: #6b7280;
        --neutral-light: #f9fafb;
        --success-green: #10b981;
        --warning-amber: #f59e0b;
        --gradient-primary: linear-gradient(135deg, #6366f1 0%, #8b5cf6 100%);
        --gradient-secondary: linear-gradient(135deg, #14b8a6 0%, #06b6d4 100%);
        --gradient-warm: linear-gradient(135deg, #f97316 0%, #ec4899 100%);
        --gradient-sophisticated: linear-gradient(135deg, #8b5cf6 0%, #a78bfa 100%);
        --shadow-elegant: 0 20px 40px rgba(0, 0, 0, 0.12);
        --shadow-soft: 0 4px 20px rgba(0, 0, 0, 0.08);
    }
    
    /* Enhanced typography */
    .main-title {
        font-family: 'Space Grotesk', sans-serif;
        font-size: 3.8rem;
        font-weight: 700;
        background: var(--gradient-primary);
        -webkit-background-clip: text;
        -webkit-text-fill-color: transparent;
        background-clip: text;
        text-align: center;
        margin-bottom: 0.5rem;
        line-height: 1.2;
        letter-spacing: -0.02em;
    }
    
    .subtitle {
        font-family: 'Inter', sans-serif;
        font-size: 1.4rem;
        color: var(--neutral-medium);
        text-align: center;
        margin-bottom: 2rem;
        font-weight: 400;
    }
    
    /* Enhanced hero section */
    .hero-container {
        background: var(--gradient-sophisticated);
        border-radius: 28px;
        padding: 3.5rem 2.5rem;
        margin: 2rem 0;
        text-align: center;
        position: relative;
        overflow: hidden;
        box-shadow: var(--shadow-elegant);
    }
    
    .hero-container::before {
        content: '';
        position: absolute;
        top: 0;
        left: 0;
        right: 0;
        bottom: 0;
        background: url('data:image/svg+xml,<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 100 100"><defs><pattern id="grain" width="100" height="100" patternUnits="userSpaceOnUse"><circle cx="50" cy="50" r="1" fill="white" opacity="0.1"/></pattern></defs><rect width="100" height="100" fill="url(%23grain)"/></svg>');
        opacity: 0.3;
    }
    
    .hero-title {
        font-family: 'Space Grotesk', sans-serif;
        font-size: 2.8rem;
        font-weight: 700;
        color: white;
        margin-bottom: 1rem;
        position: relative;
        z-index: 1;
        letter-spacing: -0.01em;
    }
    
    .hero-subtitle {
        font-family: 'Inter', sans-serif;
        font-size: 1.2rem;
        color: rgba(255, 255, 255, 0.9);
        margin-bottom: 0;
        position: relative;
        z-index: 1;
    }
    
    /* Enhanced card containers */
    .elegant-card {
        background: white;
        border-radius: 24px;
        padding: 2.5rem;
        margin: 1.5rem 0;
        box-shadow: var(--shadow-soft);
        border: 1px solid rgba(99, 102, 241, 0.08);
        transition: all 0.4s ease;
        position: relative;
        overflow: hidden;
    }
    
    .elegant-card::before {
        content: '';
        position: absolute;
        top: 0;
        left: 0;
        width: 100%;
        height: 4px;
        background: var(--gradient-sophisticated);
    }
    
    .elegant-card:hover {
        transform: translateY(-4px);
        box-shadow: var(--shadow-elegant);
    }
    
    /* Advanced theme cards */
    .theme-card {
        background: white;
        border-radius: 20px;
        padding: 2rem;
        margin: 1rem 0;
        box-shadow: var(--shadow-soft);
        border: 2px solid #e5e7eb;
        transition: all 0.4s ease;
        cursor: pointer;
        text-align: center;
        position: relative;
        overflow: hidden;
    }
    
    .theme-card:hover {
        transform: translateY(-6px);
        box-shadow: var(--shadow-elegant);
        border-color: var(--primary-purple);
    }
    
    .theme-card.selected {
        border-color: var(--primary-purple);
        background: linear-gradient(135deg, #eff6ff, #dbeafe);
        transform: translateY(-6px);
        box-shadow: 0 16px 40px rgba(99, 102, 241, 0.25);
    }
    
    .theme-preview {
        width: 100%;
        height: 140px;
        border-radius: 16px;
        margin-bottom: 1.5rem;
        display: flex;
        align-items: center;
        justify-content: center;
        color: white;
        font-weight: 600;
        font-size: 1.2rem;
        position: relative;
        overflow: hidden;
    }
    
    .theme-minimalist {
        background: linear-gradient(135deg, #f8fafc, #e2e8f0);
        color: #1e293b;
        border: 2px solid #e2e8f0;
    }
    
    .theme-dark {
        background: linear-gradient(135deg, #1e293b, #0f172a);
        color: white;
    }
    
    .theme-colorful {
        background: linear-gradient(135deg, #ff8c00, #ff66b2);
        color: white;
    }
    
    .theme-professional {
        background: linear-gradient(135deg, #1f2937, #374151);
        color: white;
    }
    
    .theme-academic {
        background: linear-gradient(135deg, #065f46, #047857);
        color: white;
    }
    
    .theme-modern {
        background: linear-gradient(135deg, #7c3aed, #a855f7);
        color: white;
    }
    
    /* Enhanced status boxes */
    .status-success {
        background: linear-gradient(135deg, #ecfdf5 0%, #d1fae5 100%);
        border: 1px solid #a7f3d0;
        border-radius: 20px;
        padding: 2rem;
        margin: 1rem 0;
        color: #065f46;
        font-family: 'Inter', sans-serif;
        font-weight: 500;
        box-shadow: var(--shadow-soft);
    }
    
    .status-info {
        background: linear-gradient(135deg, #eff6ff 0%, #dbeafe 100%);
        border: 1px solid #93c5fd;
        border-radius: 20px;
        padding: 2rem;
        margin: 1rem 0;
        color: #1e3a8a;
        font-family: 'Inter', sans-serif;
        font-weight: 500;
        box-shadow: var(--shadow-soft);
    }
    
    .status-warning {
        background: linear-gradient(135deg, #fffbeb 0%, #fef3c7 100%);
        border: 1px solid #fcd34d;
        border-radius: 20px;
        padding: 2rem;
        margin: 1rem 0;
        color: #92400e;
        font-family: 'Inter', sans-serif;
        font-weight: 500;
        box-shadow: var(--shadow-soft);
    }
    
    .status-research {
        background: linear-gradient(135deg, #fdf4ff 0%, #fae8ff 100%);
        border: 1px solid #d8b4fe;
        border-radius: 20px;
        padding: 2rem;
        margin: 1rem 0;
        color: #7c2d12;
        font-family: 'Inter', sans-serif;
        font-weight: 500;
        box-shadow: var(--shadow-soft);
    }
    
    /* Enhanced progress tracking */
    .progress-container {
        background: white;
        border-radius: 20px;
        padding: 2rem;
        margin: 1rem 0;
        box-shadow: var(--shadow-soft);
        border: 1px solid #f3f4f6;
    }
    
    .progress-step {
        display: flex;
        align-items: center;
        padding: 1rem 0;
        font-family: 'Inter', sans-serif;
        font-weight: 500;
        border-radius: 12px;
        margin: 0.5rem 0;
        transition: all 0.3s ease;
        position: relative;
    }
    
    .progress-step.completed {
        background: linear-gradient(135deg, #ecfdf5, #d1fae5);
        color: #065f46;
        transform: scale(1.02);
    }
    
    .progress-step.active {
        background: linear-gradient(135deg, #eff6ff, #dbeafe);
        color: #1e3a8a;
        transform: scale(1.05);
        box-shadow: 0 4px 12px rgba(99, 102, 241, 0.2);
    }
    
    .progress-step.pending {
        background: #f9fafb;
        color: #6b7280;
    }
    
    /* Advanced feature sections */
    .feature-grid {
        display: grid;
        grid-template-columns: repeat(auto-fit, minmax(280px, 1fr));
        gap: 1.5rem;
        margin: 2rem 0;
    }
    
    .feature-card {
        background: white;
        border-radius: 16px;
        padding: 1.5rem;
        box-shadow: var(--shadow-soft);
        border: 1px solid #f3f4f6;
        transition: all 0.3s ease;
    }
    
    .feature-card:hover {
        transform: translateY(-2px);
        box-shadow: var(--shadow-elegant);
    }
    
    .feature-icon {
        width: 48px;
        height: 48px;
        border-radius: 12px;
        background: var(--gradient-sophisticated);
        display: flex;
        align-items: center;
        justify-content: center;
        color: white;
        font-size: 1.5rem;
        margin-bottom: 1rem;
    }
    
    /* Research preview styling */
    .research-preview {
        background: linear-gradient(135deg, #fdf4ff, #fae8ff);
        border-radius: 20px;
        padding: 2rem;
        margin: 1.5rem 0;
        border-left: 4px solid var(--primary-purple);
        box-shadow: var(--shadow-soft);
    }
    
    /* Slide preview enhancements */
    .slide-preview {
        background: white;
        border-radius: 20px;
        padding: 2rem;
        margin: 1.5rem 0;
        box-shadow: var(--shadow-soft);
        border: 1px solid #e5e7eb;
        transition: all 0.3s ease;
    }
    
    .slide-preview:hover {
        transform: translateY(-2px);
        box-shadow: var(--shadow-elegant);
    }
    
    /* Enhanced metrics */
    .metric-card {
        background: white;
        border-radius: 20px;
        padding: 2rem;
        text-align: center;
        box-shadow: var(--shadow-soft);
        border: 1px solid #e5e7eb;
        transition: all 0.3s ease;
        position: relative;
        overflow: hidden;
    }
    
    .metric-card::before {
        content: '';
        position: absolute;
        top: 0;
        left: 0;
        width: 100%;
        height: 3px;
        background: var(--gradient-sophisticated);
    }
    
    .metric-card:hover {
        transform: translateY(-4px);
        box-shadow: var(--shadow-elegant);
    }
    
    .metric-value {
        font-family: 'Space Grotesk', sans-serif;
        font-size: 2.5rem;
        font-weight: 700;
        color: var(--primary-purple);
        margin-bottom: 0.5rem;
    }
    
    .metric-label {
        font-family: 'Inter', sans-serif;
        font-size: 0.9rem;
        color: var(--neutral-medium);
        font-weight: 500;
        text-transform: uppercase;
        letter-spacing: 0.05em;
    }
    
    /* Enhanced button styling */
    .stButton > button {
        border-radius: 16px !important;
        font-family: 'Inter', sans-serif !important;
        font-weight: 600 !important;
        padding: 1rem 2.5rem !important;
        transition: all 0.3s ease !important;
        border: none !important;
        text-transform: none !important;
        letter-spacing: 0.01em !important;
    }
    
    .stButton > button[kind="primary"] {
        background: var(--gradient-sophisticated) !important;
        color: white !important;
        box-shadow: 0 4px 12px rgba(139, 92, 246, 0.3) !important;
    }
    
    .stButton > button[kind="primary"]:hover {
        transform: translateY(-2px) !important;
        box-shadow: 0 12px 28px rgba(139, 92, 246, 0.4) !important;
    }
    
    /* Advanced loading animations */
    @keyframes sophisticatedPulse {
        0%, 100% { 
            opacity: 1; 
            transform: scale(1);
        }
        50% { 
            opacity: 0.8; 
            transform: scale(1.02);
        }
    }
    
    .loading-pulse {
        animation: sophisticatedPulse 2.5s infinite;
    }
    
    /* Advanced tooltips */
    .tooltip {
        position: relative;
        cursor: help;
    }
    
    .tooltip:hover::after {
        content: attr(data-tooltip);
        position: absolute;
        bottom: 100%;
        left: 50%;
        transform: translateX(-50%);
        background: var(--neutral-dark);
        color: white;
        padding: 0.5rem 1rem;
        border-radius: 8px;
        font-size: 0.8rem;
        white-space: nowrap;
        z-index: 1000;
    }
    
    /* Responsive enhancements */
    @media (max-width: 768px) {
        .main-title {
            font-size: 2.8rem;
        }
        
        .hero-title {
            font-size: 2.2rem;
        }
        
        .elegant-card {
            padding: 1.5rem;
            margin: 1rem 0;
        }
        
        .theme-card {
            margin: 0.5rem 0;
            padding: 1.5rem;
        }
        
        .feature-grid {
            grid-template-columns: 1fr;
        }
    }
</style>
""", unsafe_allow_html=True)

# Enhanced session state initialization
if 'lesson_data' not in st.session_state:
    st.session_state.lesson_data = {}
if 'current_step' not in st.session_state:
    st.session_state.current_step = 1
if 'generated_content' not in st.session_state:
    st.session_state.generated_content = None
if 'slides_approved' not in st.session_state:
    st.session_state.slides_approved = False
if 'selected_theme' not in st.session_state:
    st.session_state.selected_theme = "minimalist"
if 'research_data' not in st.session_state:
    st.session_state.research_data = {}
if 'advanced_features' not in st.session_state:
    st.session_state.advanced_features = {}
if 'generation_history' not in st.session_state:
    st.session_state.generation_history = []

class EnhancedLessonGenerator:
    def __init__(self, claude_key: str, elevenlabs_key: str):
        self.claude_key = claude_key
        self.elevenlabs_key = elevenlabs_key
        self.client = anthropic.Anthropic(api_key=claude_key)
        
    def extract_text_from_file(self, uploaded_file) -> str:
        """Enhanced text extraction with better error handling"""
        try:
            if uploaded_file.type == "text/plain":
                content = str(uploaded_file.read(), "utf-8")
                # Clean and normalize content
                content = re.sub(r'\s+', ' ', content).strip()
                return content
            else:
                return "Please use TXT files for best compatibility."
        except Exception as e:
            return f"Error reading file: {str(e)}"
    
    def conduct_deep_research(self, topic: str, content: str) -> Dict:
        """Conduct comprehensive research on the topic"""
        try:
            research_prompt = f"""
            Conduct comprehensive research on "{topic}" using the provided content as a foundation:
            
            Content: {content[:2000]}
            
            Provide detailed analysis in these areas:
            1. **Historical Context**: Timeline, key events, evolution
            2. **Current Relevance**: Modern applications, recent developments
            3. **Key Figures**: Important people, their contributions
            4. **Technical Details**: Specifications, processes, mechanisms
            5. **Cultural Impact**: Social significance, global influence
            6. **Future Outlook**: Trends, predictions, potential developments
            7. **Interesting Connections**: Related topics, surprising links
            8. **Visual Elements**: Specific image suggestions for presentation
            
            Make this research comprehensive and presentation-ready.
            """
            
            response = self.client.messages.create(
                model="claude-3-5-sonnet-20241022",
                max_tokens=2500,
                temperature=0.7,
                messages=[{"role": "user", "content": research_prompt}]
            )
            
            return {
                "content": response.content[0].text,
                "timestamp": datetime.now().isoformat(),
                "topic": topic,
                "depth": "comprehensive"
            }
        except Exception as e:
            st.error(f"Research error: {str(e)}")
            return {"content": f"Research failed for {topic}", "topic": topic}
    
    def get_interesting_facts(self, topic: str, content: str, research_data: Dict = None) -> str:
        """Enhanced fact generation with research integration"""
        try:
            research_context = ""
            if research_data and 'content' in research_data:
                research_context = f"Additional research context: {research_data['content'][:1000]}"
            
            prompt = f"""Based on the topic "{topic}", the following content, and additional research, find 7-10 fascinating and engaging facts:

Content: {content[:2000]}
{research_context}

Focus on:
- Surprising statistics and data points
- Historical anecdotes and stories
- Real-world applications and impact
- Fun trivia and lesser-known facts
- Current relevance and modern connections
- Technical innovations and breakthroughs
- Cultural and social significance

Format as a numbered list with detailed explanations that will captivate students."""

            response = self.client.messages.create(
                model="claude-3-5-sonnet-20241022",
                max_tokens=1200,
                temperature=0.7,
                messages=[{"role": "user", "content": prompt}]
            )
            return response.content[0].text
        except Exception as e:
            st.error(f"Error generating facts: {str(e)}")
            return f"Unable to generate facts due to API error."
    
    def create_advanced_lesson_outline(self, objectives: str, content: str, facts: str, research_data: Dict = None) -> str:
        """Create a sophisticated lesson outline with research integration"""
        try:
            research_context = ""
            if research_data and 'content' in research_data:
                research_context = f"Research insights: {research_data['content'][:1500]}"
            
            prompt = f"""Create a sophisticated, research-driven lesson outline:

Learning Objectives: {objectives}
Content Material: {content[:1500]}
Interesting Facts: {facts}
{research_context}

Structure the lesson with:
1. **Hook/Introduction** (5-8 minutes) - Engaging opener
2. **Historical Foundation** (10-15 minutes) - Background and context
3. **Core Concepts** (15-20 minutes) - Main content with examples
4. **Modern Relevance** (8-12 minutes) - Current applications
5. **Interactive Exploration** (10-15 minutes) - Activities and discussion
6. **Synthesis & Reflection** (5-10 minutes) - Wrap-up and assessment

For each section, include:
- Key talking points
- Suggested activities or interactions
- Visual aids needed
- Transition strategies
- Assessment opportunities

Make this outline comprehensive and engaging for modern learners."""

            response = self.client.messages.create(
                model="claude-3-5-sonnet-20241022",
                max_tokens=2000,
                temperature=0.6,
                messages=[{"role": "user", "content": prompt}]
            )
            return response.content[0].text
        except Exception as e:
            st.error(f"Error creating outline: {str(e)}")
            return f"Unable to generate lesson outline due to API error."
    
    def generate_enhanced_slide_content(self, outline: str, objectives: str, research_data: Dict = None, slide_count: int = 8) -> List[Dict]:
        """Generate sophisticated slide content with research integration"""
        try:
            research_context = ""
            if research_data and 'content' in research_data:
                research_context = f"Research insights: {research_data['content'][:1000]}"
            
            prompt = f"""Create content for {slide_count} sophisticated presentation slides:

Lesson Outline: {outline}
Objectives: {objectives}
{research_context}

For each slide, provide:
1. Slide title (engaging and descriptive)
2. Subtitle (optional, for context)
3. Key bullet points (3-4 impactful points)
4. Speaker notes (detailed, 3-4 sentences)
5. Image description (specific visual suggestions)
6. Layout style (title_slide, content_image, image_focus, split_content, full_image)
7. Design notes (specific styling suggestions)

Return ONLY valid JSON in this exact format:
[
    {{
        "slide_number": 1,
        "title": "Compelling Slide Title",
        "subtitle": "Optional descriptive subtitle",
        "content": ["Impactful point 1", "Engaging point 2", "Memorable point 3"],
        "speaker_notes": "Detailed explanation with context and examples...",
        "image_description": "Specific, detailed image suggestion",
        "layout_style": "title_slide",
        "design_notes": "Specific visual design guidance"
    }}
]

Make slides visually engaging and content-rich."""

            response = self.client.messages.create(
                model="claude-3-5-sonnet-20241022",
                max_tokens=2500,
                temperature=0.6,
                messages=[{"role": "user", "content": prompt}]
            )
            
            # Parse JSON response
            content = response.content[0].text.strip()
            # Remove any markdown formatting
            if content.startswith("json"):
                content = content[7:]
            if content.endswith(""):
                content = content[:-3]
            
            slides_content = json.loads(content)
            return slides_content
        except json.JSONDecodeError as e:
            st.error(f"Error parsing slide content JSON: {str(e)}")
            return self._get_enhanced_fallback_slides(slide_count)
        except Exception as e:
            st.error(f"Error generating slides: {str(e)}")
            return self._get_enhanced_fallback_slides(slide_count)
    
    def _get_enhanced_fallback_slides(self, count: int = 8) -> List[Dict]:
        """Enhanced fallback slide structure"""
        fallback_slides = [
            {
                "slide_number": 1,
                "title": "Introduction & Overview",
                "subtitle": "Setting the foundation",
                "content": ["Welcome and context", "Learning objectives", "Journey ahead"],
                "speaker_notes": "Welcome students and establish the learning context. Clearly outline what they will discover and why it matters.",
                "image_description": "Engaging hero image representing the topic",
                "layout_style": "title_slide",
                "design_notes": "Bold, welcoming design with strong visual impact"
            },
            {
                "slide_number": 2,
                "title": "Historical Foundation",
                "subtitle": "Understanding the origins",
                "content": ["Origins and background", "Key historical moments", "Evolution over time"],
                "speaker_notes": "Provide essential historical context that helps students understand how we arrived at current understanding.",
                "image_description": "Historical photograph or timeline visualization",
                "layout_style": "content_image",
                "design_notes": "Balanced layout with historical imagery"
            }
        ]
        
        # Add more slides to reach desired count
        for i in range(3, count + 1):
            fallback_slides.append({
                "slide_number": i,
                "title": f"Key Concept {i-1}",
                "subtitle": "Building understanding",
                "content": [f"Important aspect {i-1}", "Supporting details", "Real-world connections"],
                "speaker_notes": f"Explore this key concept with examples and connections to student experience.",
                "image_description": f"Relevant illustration for concept {i-1}",
                "layout_style": "content_image",
                "design_notes": "Clear, focused design emphasizing key points"
            })
        
        return fallback_slides[:count]
    
    def apply_enhanced_themes(self, prs, slide, theme: str, slide_data: Dict = None):
        """Apply enhanced theme styling with sophisticated design"""
        
        # Theme configurations
        themes = {
            "minimalist": {
                "bg_color": RGBColor(255, 255, 255),
                "title_font": "Segoe UI",
                "title_size": Pt(36),
                "title_color": RGBColor(30, 41, 55),
                "content_font": "Segoe UI",
                "content_size": Pt(18),
                "content_color": RGBColor(55, 65, 81)
            },
            "dark": {
                "bg_color": RGBColor(15, 23, 42),
                "title_font": "Calibri",
                "title_size": Pt(36),
                "title_color": RGBColor(248, 250, 252),
                "content_font": "Calibri",
                "content_size": Pt(18),
                "content_color": RGBColor(226, 232, 240)
            },
            "colorful": {
                "bg_color": RGBColor(255, 102, 178),
                "title_font": "Arial Rounded MT Bold",
                "title_size": Pt(36),
                "title_color": RGBColor(255, 255, 255),
                "content_font": "Arial Rounded MT Bold",
                "content_size": Pt(18),
                "content_color": RGBColor(255, 255, 255)
            },
            "professional": {
                "bg_color": RGBColor(248, 250, 252),
                "title_font": "Calibri",
                "title_size": Pt(40),
                "title_color": RGBColor(31, 41, 55),
                "content_font": "Calibri",
                "content_size": Pt(18),
                "content_color": RGBColor(75, 85, 99)
            },
            "academic": {
                "bg_color": RGBColor(245, 248, 250),
                "title_font": "Times New Roman",
                "title_size": Pt(38),
                "title_color": RGBColor(6, 95, 70),
                "content_font": "Times New Roman",
                "content_size": Pt(16),
                "content_color": RGBColor(17, 24, 39)
            },
            "modern": {
                "bg_color": RGBColor(250, 248, 255),
                "title_font": "Space Grotesk",
                "title_size": Pt(38),
                "title_color": RGBColor(124, 58, 237),
                "content_font": "Inter",
                "content_size": Pt(17),
                "content_color": RGBColor(55, 65, 81)
            }
        }
        
        theme_config = themes.get(theme, themes["minimalist"])
        
        # Apply background
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = theme_config["bg_color"]
        
        # Style title
        if hasattr(slide.shapes, 'title') and slide.shapes.title:
            title_shape = slide.shapes.title
            if title_shape.text_frame.paragraphs:
                title_para = title_shape.text_frame.paragraphs[0]
                title_para.font.name = theme_config["title_font"]
                title_para.font.size = theme_config["title_size"]
                title_para.font.color.rgb = theme_config["title_color"]
                title_para.font.bold = True
        
        # Style content
        for shape in slide.shapes:
            if hasattr(shape, 'text_frame') and shape != slide.shapes.title:
                for paragraph in shape.text_frame.paragraphs:
                    paragraph.font.name = theme_config["content_font"]
                    paragraph.font.size = theme_config["content_size"]
                    paragraph.font.color.rgb = theme_config["content_color"]
    
    def create_sophisticated_powerpoint(self, slides_data: List[Dict], lesson_title: str, theme: str = "minimalist") -> io.BytesIO:
        """Create sophisticated PowerPoint with enhanced themes"""
        try:
            if not slides_data or not isinstance(slides_data, list):
                st.error("Invalid slide data provided")
                return None
                
            prs = Presentation()
            
            # Process each slide with enhanced styling
            for slide_data in slides_data:
                try:
                    layout_style = slide_data.get('layout_style', 'content_image')
                    
                    if layout_style == 'title_slide':
                        slide_layout = prs.slide_layouts[0]  # Title slide
                    else:
                        slide_layout = prs.slide_layouts[1]  # Content slide
                    
                    slide = prs.slides.add_slide(slide_layout)
                    
                    # Add title
                    if hasattr(slide.shapes, 'title') and slide.shapes.title:
                        slide.shapes.title.text = slide_data.get('title', '')
                    
                    # Add subtitle for title slides
                    if layout_style == 'title_slide' and len(slide.placeholders) > 1:
                        subtitle = slide.placeholders[1]
                        subtitle_text = slide_data.get('subtitle', 'AI-Generated Educational Content')
                        subtitle.text = subtitle_text
                    
                    # Add content for content slides
                    elif layout_style != 'title_slide':
                        if len(slide.placeholders) > 1:
                            content_placeholder = slide.placeholders[1]
                            text_frame = content_placeholder.text_frame
                            text_frame.clear()
                            
                            # Add subtitle if present
                            if slide_data.get('subtitle'):
                                p = text_frame.add_paragraph()
                                p.text = slide_data['subtitle']
                                p.level = 0
                                p.font.size = Pt(20)
                                p.font.italic = True
                                
                            # Add content points
                            for point in slide_data.get('content', []):
                                if point and isinstance(point, str):
                                    p = text_frame.add_paragraph()
                                    p.text = str(point)
                                    p.level = 0
                    
                    # Apply enhanced theme
                    self.apply_enhanced_themes(prs, slide, theme, slide_data)
                        
                except Exception as slide_error:
                    st.warning(f"Error creating slide {slide_data.get('slide_number', 'unknown')}: {str(slide_error)}")
                    continue
            
            # Save to BytesIO
            pptx_buffer = io.BytesIO()
            prs.save(pptx_buffer)
            pptx_buffer.seek(0)
            
            return pptx_buffer
        except Exception as e:
            st.error(f"Error creating PowerPoint: {str(e)}")
            return None
    
    def generate_enhanced_audio(self, text: str, voice_id: str = "21m00Tcm4TlvDq8ikWAM", 
                               stability: float = 0.5, similarity: float = 0.5) -> bytes:
        """Enhanced audio generation with customizable parameters"""
        try:
            url = f"https://api.elevenlabs.io/v1/text-to-speech/{voice_id}"
            
            headers = {
                "Accept": "audio/mpeg",
                "Content-Type": "application/json",
                "xi-api-key": self.elevenlabs_key
            }
            
            data = {
                "text": text,
                "model_id": "eleven_monolingual_v1",
                "voice_settings": {
                    "stability": stability,
                    "similarity_boost": similarity
                }
            }
            
            response = requests.post(url, json=data, headers=headers)
            
            if response.status_code == 200:
                return response.content
            else:
                st.error(f"ElevenLabs API error: {response.status_code}")
                return None
        except Exception as e:
            st.error(f"Error generating audio: {str(e)}")
            return None

def render_enhanced_theme_selector():
    """Enhanced theme selection with more options"""
    st.markdown("<h3 style='color: #6366f1; font-family: Inter, sans-serif; margin: 2rem 0 1rem 0;'>🎨 Choose Presentation Theme</h3>", unsafe_allow_html=True)
    
    themes = [
        {
            "key": "minimalist",
            "name": "Minimalist",
            "description": "Clean & Simple",
            "preview_class": "theme-minimalist",
            "best_for": "Professional presentations"
        },
        {
            "key": "dark",
            "name": "Dark Mode",
            "description": "Elegant & Bold",
            "preview_class": "theme-dark",
            "best_for": "Tech & Modern topics"
        },
        {
            "key": "colorful",
            "name": "Colorful",
            "description": "Vibrant & Creative",
            "preview_class": "theme-colorful",
            "best_for": "Creative & Arts subjects"
        },
        {
            "key": "professional",
            "name": "Professional",
            "description": "Corporate & Formal",
            "preview_class": "theme-professional",
            "best_for": "Business presentations"
        },
        {
            "key": "academic",
            "name": "Academic",
            "description": "Scholarly & Traditional",
            "preview_class": "theme-academic",
            "best_for": "Educational institutions"
        },
        {
            "key": "modern",
            "name": "Modern",
            "description": "Contemporary & Sleek",
            "preview_class": "theme-modern",
            "best_for": "Innovation & Technology"
        }
    ]
    
    # Create theme grid
    cols = st.columns(3)
    
    for i, theme in enumerate(themes):
        with cols[i % 3]:
            selected_class = "selected" if st.session_state.selected_theme == theme["key"] else ""
            
            theme_html = f"""
            <div class="theme-card {selected_class}">
                <div class="theme-preview {theme['preview_class']}">
                    {theme['name']}
                </div>
                <h4 style="color: #1f2937; margin: 0.5rem 0; font-family: 'Inter', sans-serif;">{theme['name']}</h4>
                <p style="color: #6b7280; margin: 0.5rem 0; font-size: 0.875rem;">{theme['description']}</p>
                <small style="color: #9ca3af; font-size: 0.75rem;">Best for: {theme['best_for']}</small>
            </div>
            """
            
            st.markdown(theme_html, unsafe_allow_html=True)
            
            if st.button(f"Select {theme['name']}", key=f"theme_{theme['key']}", use_container_width=True):
                st.session_state.selected_theme = theme["key"]
                st.rerun()
    
    # Show selected theme with details
    selected_theme = next(t for t in themes if t["key"] == st.session_state.selected_theme)
    st.markdown(f"""
    <div class="status-info" style="text-align: center; margin-top: 1.5rem;">
        🎨 <strong>Selected Theme:</strong> {selected_theme['name']} - {selected_theme['description']}<br>
        <small>Perfect for: {selected_theme['best_for']}</small>
    </div>
    """, unsafe_allow_html=True)

def render_advanced_features():
    """Render advanced feature configuration"""
    st.markdown("## 🚀 Advanced Features")
    
    with st.expander("🔬 Research Enhancement", expanded=False):
        col1, col2 = st.columns(2)
        
        with col1:
            st.markdown("### Research Depth")
            research_depth = st.selectbox(
                "Research Level",
                ["Standard", "Comprehensive", "Deep Dive"],
                index=1,
                help="Higher levels provide more detailed research but take longer"
            )
            
            auto_research = st.checkbox(
                "Auto-Research Topics",
                value=True,
                help="Automatically conduct research on unknown topics"
            )
            
        with col2:
            st.markdown("### Content Focus")
            focus_areas = st.multiselect(
                "Research Focus Areas",
                ["Historical Context", "Technical Details", "Cultural Impact", 
                 "Current Relevance", "Future Trends", "Global Perspective"],
                default=["Historical Context", "Current Relevance"],
                help="Select areas to emphasize in research"
            )
    
    with st.expander("🎨 Visual Customization", expanded=False):
        col1, col2, col3 = st.columns(3)
        
        with col1:
            st.markdown("### Slide Design")
            slide_count = st.slider(
                "Target Slide Count",
                min_value=4,
                max_value=15,
                value=8,
                help="Recommended: 6-10 slides for optimal engagement"
            )
            
            layout_preference = st.selectbox(
                "Layout Preference",
                ["Balanced", "Text-Heavy", "Visual-Heavy", "Minimalist"],
                help="Influences text-to-image ratio"
            )
        
        with col2:
            st.markdown("### Typography")
            font_style = st.selectbox(
                "Font Style",
                ["Modern", "Classic", "Academic", "Creative"],
                help="Affects font choices across themes"
            )
            
            text_density = st.selectbox(
                "Text Density",
                ["Minimal", "Moderate", "Detailed"],
                index=1,
                help="Amount of text per slide"
            )
        
        with col3:
            st.markdown("### Color & Mood")
            color_intensity = st.slider(
                "Color Intensity",
                min_value=0.3,
                max_value=1.0,
                value=0.7,
                step=0.1,
                help="Adjust overall color saturation"
            )
            
            mood = st.selectbox(
                "Presentation Mood",
                ["Professional", "Energetic", "Calm", "Inspiring"],
                help="Influences design choices"
            )
    
    with st.expander("🎙️ Audio Enhancement", expanded=False):
        col1, col2 = st.columns(2)
        
        with col1:
            st.markdown("### Voice Settings")
            voice_style = st.selectbox(
                "Voice Style",
                ["Professional", "Conversational", "Energetic", "Calm"],
                help="Affects stability and similarity settings"
            )
            
            speaking_pace = st.selectbox(
                "Speaking Pace",
                ["Slow", "Normal", "Fast"],
                index=1,
                help="Adjusts text processing for audio"
            )
        
        with col2:
            st.markdown("### Audio Quality")
            audio_format = st.selectbox(
                "Audio Format",
                ["MP3 (Standard)", "MP3 (High Quality)"],
                help="Higher quality = larger file sizes"
            )
            
            include_pauses = st.checkbox(
                "Smart Pauses",
                value=True,
                help="Add natural pauses between points"
            )
    
    # Store advanced settings
    st.session_state.advanced_features = {
        "research_depth": locals().get('research_depth', 'Comprehensive'),
        "auto_research": locals().get('auto_research', True),
        "focus_areas": locals().get('focus_areas', []),
        "slide_count": locals().get('slide_count', 8),
        "layout_preference": locals().get('layout_preference', 'Balanced'),
        "font_style": locals().get('font_style', 'Modern'),
        "text_density": locals().get('text_density', 'Moderate'),
        "color_intensity": locals().get('color_intensity', 0.7),
        "mood": locals().get('mood', 'Professional'),
        "voice_style": locals().get('voice_style', 'Professional'),
        "speaking_pace": locals().get('speaking_pace', 'Normal'),
        "audio_format": locals().get('audio_format', 'MP3 (Standard)'),
        "include_pauses": locals().get('include_pauses', True)
    }

def render_generation_analytics():
    """Show analytics and insights about the generation process"""
    if st.session_state.generation_history:
        st.markdown("## 📊 Generation Analytics")
        
        col1, col2, col3, col4 = st.columns(4)
        
        with col1:
            total_generations = len(st.session_state.generation_history)
            st.metric("Total Generations", total_generations)
        
        with col2:
            avg_slides = np.mean([g.get('slide_count', 0) for g in st.session_state.generation_history])
            st.metric("Avg. Slides", f"{avg_slides:.1f}")
        
        with col3:
            themes_used = set([g.get('theme', 'unknown') for g in st.session_state.generation_history])
            st.metric("Themes Used", len(themes_used))
        
        with col4:
            success_rate = len([g for g in st.session_state.generation_history if g.get('success', False)]) / total_generations * 100
            st.metric("Success Rate", f"{success_rate:.0f}%")

def render_smart_suggestions():
    """Provide intelligent suggestions based on content"""
    if 'lesson_data' in st.session_state and st.session_state.lesson_data:
        data = st.session_state.lesson_data
        
        st.markdown("## 💡 Smart Suggestions")
        
        suggestions = []
        
        # Analyze content for suggestions
        if 'content' in data:
            content = data['content'].lower()
            
            if any(word in content for word in ['history', 'historical', 'ancient', 'century']):
                suggestions.append({
                    "icon": "🏛️",
                    "title": "Historical Timeline",
                    "description": "Consider adding a timeline slide to show chronological progression"
                })
            
            if any(word in content for word in ['process', 'steps', 'procedure', 'method']):
                suggestions.append({
                    "icon": "🔄",
                    "title": "Process Diagram",
                    "description": "A flowchart or step-by-step visual would enhance understanding"
                })
            
            if any(word in content for word in ['data', 'statistics', 'numbers', 'research']):
                suggestions.append({
                    "icon": "📈",
                    "title": "Data Visualization",
                    "description": "Consider adding charts or graphs to illustrate key statistics"
                })
            
            if any(word in content for word in ['geography', 'location', 'place', 'region']):
                suggestions.append({
                    "icon": "🗺️",
                    "title": "Map Integration",
                    "description": "Maps could help students visualize geographical concepts"
                })
        
        if suggestions:
            for suggestion in suggestions:
                st.markdown(f"""
                <div class="feature-card">
                    <div class="feature-icon">{suggestion['icon']}</div>
                    <h4>{suggestion['title']}</h4>
                    <p>{suggestion['description']}</p>
                </div>
                """, unsafe_allow_html=True)

def main():
    # Enhanced Header
    st.markdown("""
    <div class="hero-container">
        <h1 class="hero-title">🚀 Enhanced AI Lesson Generator</h1>
        <p class="hero-subtitle">Create sophisticated, research-driven presentations with advanced AI capabilities</p>
    </div>
    """, unsafe_allow_html=True)
    
    # Enhanced Sidebar
    with st.sidebar:
        st.markdown("""
        <div style="text-align: center; padding: 1rem 0;">
            <h2 style="color: #8b5cf6; font-family: 'Space Grotesk', sans-serif; margin-bottom: 0.5rem;">⚙️ Configuration</h2>
        </div>
        """, unsafe_allow_html=True)
        
        # API Keys section
        with st.expander("🔐 API Keys", expanded=True):
            claude_key = st.text_input(
                "Anthropic Claude API Key", 
                type="password", 
                help="Get from: https://console.anthropic.com/"
            )
            elevenlabs_key = st.text_input(
                "ElevenLabs API Key", 
                type="password", 
                help="Get from: https://elevenlabs.io/"
            )
        
        if not claude_key or not elevenlabs_key:
            st.markdown("""
            <div class="status-warning">
                ⚠️ Please enter both API keys to continue
            </div>
            """, unsafe_allow_html=True)
            return
        
        # Enhanced Progress tracking
        st.markdown("<h3 style='color: #8b5cf6; font-family: Inter, sans-serif; margin: 1.5rem 0 1rem 0;'>📊 Progress</h3>", unsafe_allow_html=True)
        
        progress_steps = [
            ("📝", "Setup & Input"),
            ("🔬", "Research & Analysis"), 
            ("👀", "Review & Refine"),
            ("🎨", "Theme & Style"),
            ("🎬", "Generate Content"),
            ("🎉", "Download & Share")
        ]
        
        progress_html = '<div class="progress-container">'
        for i, (icon, step) in enumerate(progress_steps, 1):
            if i < st.session_state.current_step:
                status_class = "completed"
                status_icon = "✅"
            elif i == st.session_state.current_step:
                status_class = "active"
                status_icon = "🔄"
            else:
                status_class = "pending"
                status_icon = "⏳"
            
            progress_html += f'<div class="progress-step {status_class}">{status_icon} {icon} {step}</div>'
        progress_html += '</div>'
        
        st.markdown(progress_html, unsafe_allow_html=True)
        
        # Quick stats
        if st.session_state.generation_history:
            st.markdown("### 📈 Quick Stats")
            st.metric("Sessions", len(st.session_state.generation_history))
    
    # Initialize enhanced lesson generator
    if claude_key and elevenlabs_key:
        lesson_gen = EnhancedLessonGenerator(claude_key, elevenlabs_key)
    else:
        return
    
    # Main content with tabs
    tab1, tab2, tab3 = st.tabs(["🎓 Generate Lesson", "🚀 Advanced Features", "📊 Analytics"])
    
    with tab1:
        # Main generation workflow
        main_container = st.container()
        
        with main_container:
            # Step 1: Enhanced Input Collection
            if st.session_state.current_step == 1:
                st.markdown('<div class="elegant-card">', unsafe_allow_html=True)
                st.markdown("<h2 style='color: #8b5cf6; font-family: Space Grotesk, sans-serif; margin-bottom: 1.5rem;'>📝 Step 1: Enhanced Lesson Setup</h2>", unsafe_allow_html=True)
                
                col1, col2 = st.columns([2, 1])
                
                with col1:
                    lesson_title = st.text_input(
                        "Lesson Title", 
                        placeholder="e.g., The Revolutionary Impact of Renewable Energy",
                        help="Make it engaging and descriptive"
                    )
                    
                    col1a, col1b = st.columns(2)
                    with col1a:
                        subject = st.selectbox("Subject", [
                            "Science", "Technology", "Engineering", "Mathematics",
                            "History", "Geography", "Social Studies", "Literature",
                            "Arts", "Business", "Health", "Other"
                        ])
                    
                    with col1b:
                        grade_level = st.selectbox("Grade Level", [
                            "Elementary (K-5)", "Middle School (6-8)", 
                            "High School (9-12)", "College/University", "Professional"
                        ])
                    
                    duration = st.slider("Lesson Duration (minutes)", 15, 90, 45, 5)
                    
                    objectives = st.text_area(
                        "Learning Objectives", 
                        placeholder="What should students learn? Be specific about knowledge, skills, and understanding...",
                        height=120,
                        help="Clear objectives help generate better content"
                    )
                
                with col2:
                    st.markdown("### 🎯 Quick Start Templates")
                    
                    templates = {
                        "🔬 Scientific Discovery": {
                            "title": "Breakthrough in Quantum Computing",
                            "objectives": "Students will understand quantum computing principles and their revolutionary potential.",
                            "content": "Quantum computing represents a paradigm shift in computational power..."
                        },
                        "🏛️ Historical Event": {
                            "title": "The Fall of the Berlin Wall",
                            "objectives": "Students will analyze the causes and effects of the Berlin Wall's fall.",
                            "content": "The Berlin Wall stood as a symbol of division for 28 years..."
                        },
                        "🌱 Environmental Topic": {
                            "title": "Climate Change Solutions",
                            "objectives": "Students will evaluate various approaches to addressing climate change.",
                            "content": "Climate change presents unprecedented challenges requiring innovative solutions..."
                        }
                    }
                    
                    for template_name, template_data in templates.items():
                        if st.button(template_name, use_container_width=True):
                            lesson_title = template_data["title"]
                            objectives = template_data["objectives"]
                            # Auto-populate with template data
                            st.session_state.lesson_data = {
                                'title': template_data["title"],
                                'subject': 'Science',
                                'grade_level': 'High School (9-12)',
                                'duration': 45,
                                'objectives': template_data["objectives"],
                                'content': template_data["content"]
                            }
                            st.rerun()
                
                st.markdown("<h3 style='color: #14b8a6; font-family: Inter, sans-serif; margin: 2rem 0 1rem 0;'>📎 Upload Learning Material</h3>", unsafe_allow_html=True)
                
                uploaded_file = st.file_uploader(
                    "Upload Content File", 
                    type=['txt'], 
                    help="Upload TXT files with your lesson content"
                )
                
                # Enhanced processing options
                col3, col4 = st.columns(2)
                
                with col3:
                    auto_research = st.checkbox(
                        "Enable Auto-Research", 
                        value=True,
                        help="Automatically research your topic for additional context"
                    )
                
                with col4:
                    research_depth = st.selectbox(
                        "Research Depth",
                        ["Standard", "Comprehensive", "Deep Dive"],
                        index=1
                    )
                
                st.markdown('</div>', unsafe_allow_html=True)
                
                # Process button
                if lesson_title and objectives and (uploaded_file or st.session_state.lesson_data):
                    if st.button("🚀 Begin Enhanced Analysis", type="primary", use_container_width=True):
                        with st.spinner("✨ Processing content and conducting research..."):
                            if uploaded_file:
                                content = lesson_gen.extract_text_from_file(uploaded_file)
                            else:
                                content = st.session_state.lesson_data.get('content', '')
                            
                            # Conduct research if enabled
                            research_data = {}
                            if auto_research:
                                research_data = lesson_gen.conduct_deep_research(lesson_title, content)
                            
                            # Generate enhanced facts
                            facts = lesson_gen.get_interesting_facts(lesson_title, content, research_data)
                            
                            st.session_state.lesson_data = {
                                'title': lesson_title,
                                'subject': subject,
                                'grade_level': grade_level,
                                'duration': duration,
                                'objectives': objectives,
                                'content': content,
                                'facts': facts,
                                'research_data': research_data,
                                'auto_research': auto_research,
                                'research_depth': research_depth
                            }
                            st.session_state.current_step = 2
                            st.rerun()
            
            # Step 2: Enhanced Analysis and Research Review
            elif st.session_state.current_step == 2:
                st.markdown('<div class="elegant-card">', unsafe_allow_html=True)
                st.markdown("<h2 style='color: #8b5cf6; font-family: Space Grotesk, sans-serif; margin-bottom: 1.5rem;'>🔬 Step 2: Enhanced Analysis & Research</h2>", unsafe_allow_html=True)
                
                data = st.session_state.lesson_data
                
                # Research results
                if data.get('research_data') and data['research_data'].get('content'):
                    st.markdown("<h3 style='color: #8b5cf6; font-family: Inter, sans-serif;'>🔍 Research Results</h3>", unsafe_allow_html=True)
                    
                    with st.expander("📚 Comprehensive Research", expanded=True):
                        st.markdown(f"""
                        <div class="research-preview">
                            <h4>Deep Research on: {data['title']}</h4>
                            <p><strong>Research Depth:</strong> {data.get('research_depth', 'Standard')}</p>
                            <p><strong>Generated:</strong> {data['research_data'].get('timestamp', 'Unknown')}</p>
                            <div style="margin-top: 1rem;">
                                {data['research_data']['content'][:1200]}...
                            </div>
                        </div>
                        """, unsafe_allow_html=True)
                
                col1, col2 = st.columns(2)
                
                with col1:
                    st.markdown("<h3 style='color: #14b8a6; font-family: Inter, sans-serif;'>📚 Content Analysis</h3>", unsafe_allow_html=True)
                    
                    # Content metrics
                    content_length = len(data['content'])
                    word_count = len(data['content'].split())
                    
                    metrics_col1, metrics_col2 = st.columns(2)
                    with metrics_col1:
                        st.metric("Characters", f"{content_length:,}")
                    with metrics_col2:
                        st.metric("Words", f"{word_count:,}")
                    
                    st.text_area("Content Preview", data['content'][:600] + "...", height=200, disabled=True)
                    
                with col2:
                    st.markdown("<h3 style='color: #f97316; font-family: Inter, sans-serif;'>🎯 Enhanced Facts</h3>", unsafe_allow_html=True)
                    
                    st.markdown(f"""
                    <div style="background: linear-gradient(135deg, #fefbf3, #fef3e2); padding: 2rem; border-radius: 16px; border-left: 4px solid #f97316; max-height: 400px; overflow-y: auto;">
                        {data['facts']}
                    </div>
                    """, unsafe_allow_html=True)
                
                # Lesson overview with enhanced metrics
                st.markdown("<h3 style='color: #8b5cf6; font-family: Inter, sans-serif; margin: 2rem 0 1rem 0;'>📋 Enhanced Lesson Overview</h3>", unsafe_allow_html=True)
                
                overview_cols = st.columns(5)
                metrics = [
                    ("📖", "Title", data['title'][:20] + "..." if len(data['title']) > 20 else data['title']),
                    ("🎓", "Subject", data['subject']),
                    ("👥", "Level", data['grade_level'].split()[0]),
                    ("⏱️", "Duration", f"{data['duration']} min"),
                    ("🔬", "Research", "✅" if data.get('research_data') else "❌")
                ]
                
                for i, (icon, label, value) in enumerate(metrics):
                    with overview_cols[i]:
                        st.markdown(f"""
                        <div class="metric-card">
                            <div class="metric-value">{icon}</div>
                            <div class="metric-label">{label}</div>
                            <small style="color: #6b7280;">{value}</small>
                        </div>
                        """, unsafe_allow_html=True)
                
                # Objectives display
                st.markdown(f"""
                <div style="background: white; padding: 2rem; border-radius: 16px; border: 1px solid #e5e7eb; margin: 1.5rem 0; box-shadow: 0 4px 6px rgba(0, 0, 0, 0.05);">
                    <h4 style="color: #8b5cf6; margin-bottom: 1rem;">🎯 Learning Objectives</h4>
                    <p style="color: #6b7280; line-height: 1.6;">{data['objectives']}</p>
                </div>
                """, unsafe_allow_html=True)
                
                st.markdown('</div>', unsafe_allow_html=True)
                
                # Action buttons
                col1, col2, col3, col4 = st.columns(4)
                
                with col1:
                    if st.button("⬅️ Back to Setup", type="secondary"):
                        st.session_state.current_step = 1
                        st.rerun()
                
                with col2:
                    if st.button("🔄 Re-research", type="secondary"):
                        with st.spinner("🔍 Conducting new research..."):
                            new_research = lesson_gen.conduct_deep_research(data['title'], data['content'])
                            st.session_state.lesson_data['research_data'] = new_research
                            st.rerun()
                
                with col3:
                    if st.button("🎲 New Facts", type="secondary"):
                        with st.spinner("🎯 Generating new facts..."):
                            new_facts = lesson_gen.get_interesting_facts(
                                data['title'], 
                                data['content'], 
                                data.get('research_data')
                            )
                            st.session_state.lesson_data['facts'] = new_facts
                            st.rerun()
                
                with col4:
                    if st.button("✅ Create Outline", type="primary"):
                        with st.spinner("🎨 Creating enhanced lesson outline..."):
                            outline = lesson_gen.create_advanced_lesson_outline(
                                data['objectives'], 
                                data['content'], 
                                data['facts'],
                                data.get('research_data')
                            )
                            
                            # Get slide count from advanced features or use default
                            slide_count = st.session_state.advanced_features.get('slide_count', 8)
                            
                            slides = lesson_gen.generate_enhanced_slide_content(
                                outline, 
                                data['objectives'],
                                data.get('research_data'),
                                slide_count
                            )
                            
                            st.session_state.lesson_data['outline'] = outline
                            st.session_state.lesson_data['slides'] = slides
                            st.session_state.current_step = 3
                            st.rerun()
            
            # Step 3: Enhanced Review and Approval
            elif st.session_state.current_step == 3:
                st.markdown('<div class="elegant-card">', unsafe_allow_html=True)
                st.markdown("<h2 style='color: #8b5cf6; font-family: Space Grotesk, sans-serif; margin-bottom: 1.5rem;'>👀 Step 3: Enhanced Review & Approval</h2>", unsafe_allow_html=True)
                
                data = st.session_state.lesson_data
                
                # Enhanced lesson outline
                st.markdown("<h3 style='color: #14b8a6; font-family: Inter, sans-serif;'>📋 Sophisticated Lesson Outline</h3>", unsafe_allow_html=True)
                with st.expander("View Complete Enhanced Outline", expanded=True):
                    st.markdown(f"""
                    <div style="background: linear-gradient(135deg, #f0fdfa, #ecfdf5); padding: 2.5rem; border-radius: 20px; border-left: 4px solid #14b8a6; box-shadow: 0 4px 6px rgba(0, 0, 0, 0.05);">
                        {data['outline']}
                    </div>
                    """, unsafe_allow_html=True)
                
                # Enhanced slide previews
                st.markdown("<h3 style='color: #f97316; font-family: Inter, sans-serif; margin: 2rem 0 1rem 0;'>🖼️ Enhanced Slide Previews</h3>", unsafe_allow_html=True)
                
                if 'slides' in data and data['slides']:
                    # Slide summary metrics
                    slide_cols = st.columns(4)
                    with slide_cols[0]:
                        st.metric("Total Slides", len(data['slides']))
                    with slide_cols[1]:
                        avg_content = np.mean([len(slide.get('content', [])) for slide in data['slides']])
                        st.metric("Avg Points/Slide", f"{avg_content:.1f}")
                    with slide_cols[2]:
                        has_subtitles = sum([1 for slide in data['slides'] if slide.get('subtitle')])
                        st.metric("With Subtitles", has_subtitles)
                    with slide_cols[3]:
                        layout_types = set([slide.get('layout_style', 'standard') for slide in data['slides']])
                        st.metric("Layout Types", len(layout_types))
                    
                    # Slide previews
                    for i, slide in enumerate(data['slides']):
                        with st.expander(f"Slide {slide['slide_number']}: {slide['title']}", expanded=i == 0):
                            slide_col1, slide_col2 = st.columns([2, 1])
                            
                            with slide_col1:
                                st.markdown("**📝 Content Structure:**")
                                if slide.get('subtitle'):
                                    st.markdown(f"*Subtitle: {slide['subtitle']}*")
                                
                                for point in slide.get('content', []):
                                    st.markdown(f"• {point}")
                                
                                st.markdown(f"**🖼️ Visual Concept:** {slide.get('image_description', 'No image specified')}")
                                st.markdown(f"**📐 Layout Style:** {slide.get('layout_style', 'standard').replace('_', ' ').title()}")
                                
                            with slide_col2:
                                st.markdown("**🎤 Speaker Notes:**")
                                st.markdown(f"""
                                <div style="background: #f8fafc; padding: 1.5rem; border-radius: 12px; font-style: italic; border-left: 3px solid #8b5cf6;">
                                    {slide.get('speaker_notes', 'No speaker notes available')}
                                </div>
                                """, unsafe_allow_html=True)
                                
                                if slide.get('design_notes'):
                                    st.markdown("**🎨 Design Notes:**")
                                    st.markdown(f"""
                                    <div style="background: #fef3e2; padding: 1rem; border-radius: 8px; font-size: 0.9rem;">
                                        {slide['design_notes']}
                                    </div>
                                    """, unsafe_allow_html=True)
                
                st.markdown('</div>', unsafe_allow_html=True)
                
                # Smart suggestions
                render_smart_suggestions()
                
                # Action buttons
                col1, col2, col3, col4 = st.columns(4)
                
                with col1:
                    if st.button("⬅️ Back to Analysis", type="secondary"):
                        st.session_state.current_step = 2
                        st.rerun()
                
                with col2:
                    if st.button("🔄 Regenerate Slides", type="secondary"):
                        with st.spinner("🎨 Creating new slide variations..."):
                            slide_count = st.session_state.advanced_features.get('slide_count', 8)
                            new_slides = lesson_gen.generate_enhanced_slide_content(
                                data['outline'], 
                                data['objectives'],
                                data.get('research_data'),
                                slide_count
                            )
                            st.session_state.lesson_data['slides'] = new_slides
                            st.rerun()
                
                with col3:
                    if st.button("📝 Edit Slides", type="secondary"):
                        st.info("Individual slide editing coming soon!")
                
                with col4:
                    if st.button("✅ Approve & Style", type="primary"):
                        st.session_state.slides_approved = True
                        st.session_state.current_step = 4
                        st.rerun()
            
            # Step 4: Enhanced Theme Selection
            elif st.session_state.current_step == 4:
                st.markdown('<div class="elegant-card">', unsafe_allow_html=True)
                st.markdown("<h2 style='color: #8b5cf6; font-family: Space Grotesk, sans-serif; margin-bottom: 1.5rem;'>🎨 Step 4: Enhanced Theme & Styling</h2>", unsafe_allow_html=True)
                
                if not st.session_state.slides_approved:
                    st.error("Please approve the content first")
                    return
                
                # Enhanced theme selector
                render_enhanced_theme_selector()
                
                st.markdown('</div>', unsafe_allow_html=True)
                
                # Action buttons
                col1, col2, col3 = st.columns(3)
                
                with col1:
                    if st.button("⬅️ Back to Review", type="secondary"):
                        st.session_state.current_step = 3
                        st.rerun()
                
                with col2:
                    selected_theme = next(t['name'] for t in [
                        {"key": "minimalist", "name": "Minimalist"},
                        {"key": "dark", "name": "Dark Mode"},
                        {"key": "colorful", "name": "Colorful"},
                        {"key": "professional", "name": "Professional"},
                        {"key": "academic", "name": "Academic"},
                        {"key": "modern", "name": "Modern"}
                    ] if t["key"] == st.session_state.selected_theme)
                    
                    st.markdown(f"""
                    <div style="text-align: center; padding: 1.5rem; background: linear-gradient(135deg, #f8fafc, #f1f5f9); border-radius: 12px;">
                        <strong style="color: #8b5cf6;">Selected: {selected_theme}</strong>
                    </div>
                    """, unsafe_allow_html=True)
                
                with col3:
                    if st.button("🎬 Generate Materials", type="primary"):
                        st.session_state.current_step = 5
                        st.rerun()
            
            # Step 5: Enhanced Material Generation
            elif st.session_state.current_step == 5:
                st.markdown('<div class="elegant-card">', unsafe_allow_html=True)
                st.markdown("<h2 style='color: #8b5cf6; font-family: Space Grotesk, sans-serif; margin-bottom: 1.5rem;'>🎬 Step 5: Enhanced Material Generation</h2>", unsafe_allow_html=True)
                
                data = st.session_state.lesson_data
                
                if not st.session_state.slides_approved:
                    st.error("Please approve the content first")
                    return
                
                # Enhanced status tracking
                status_container = st.empty()
                
                # Generate PowerPoint with enhanced theme
                theme_name = next(t['name'] for t in [
                    {"key": "minimalist", "name": "Minimalist"},
                    {"key": "dark", "name": "Dark Mode"},
                    {"key": "colorful", "name": "Colorful"},
                    {"key": "professional", "name": "Professional"},
                    {"key": "academic", "name": "Academic"},
                    {"key": "modern", "name": "Modern"}
                ] if t["key"] == st.session_state.selected_theme)
                
                status_container.markdown(f"""
                <div class="loading-pulse" style="background: linear-gradient(135deg, #eff6ff, #dbeafe); padding: 2rem; border-radius: 16px; border-left: 4px solid #3b82f6;">
                    🎨 <strong>Creating sophisticated PowerPoint with {theme_name} theme...</strong><br>
                    <small>Applying advanced styling and research-driven content</small>
                </div>
                """, unsafe_allow_html=True)
                
                try:
                    pptx_buffer = lesson_gen.create_sophisticated_powerpoint(
                        data['slides'], 
                        data['title'], 
                        st.session_state.selected_theme
                    )
                except Exception as e:
                    st.error(f"Error creating PowerPoint: {str(e)}")
                    pptx_buffer = None
                
                if pptx_buffer:
                    status_container.markdown("""
                    <div class="loading-pulse" style="background: linear-gradient(135deg, #f0fdfa, #ecfdf5); padding: 2rem; border-radius: 16px; border-left: 4px solid #10b981;">
                        🎙️ <strong>Generating enhanced audio narration...</strong><br>
                        <small>Creating professional voice-over with advanced settings</small>
                    </div>
                    """, unsafe_allow_html=True)
                    
                    # Enhanced audio generation
                    audio_files = []
                    voice_settings = {
                        "Professional": {"stability": 0.7, "similarity": 0.8},
                        "Conversational": {"stability": 0.5, "similarity": 0.6},
                        "Energetic": {"stability": 0.4, "similarity": 0.7},
                        "Calm": {"stability": 0.8, "similarity": 0.9}
                    }
                    
                    voice_style = st.session_state.advanced_features.get('voice_style', 'Professional')
                    settings = voice_settings.get(voice_style, voice_settings['Professional'])
                    
                    for i, slide in enumerate(data['slides']):
                        status_container.markdown(f"""
                        <div class="loading-pulse" style="background: linear-gradient(135deg, #fefbf3, #fef3e2); padding: 2rem; border-radius: 16px; border-left: 4px solid #f59e0b;">
                            🔄 <strong>Generating audio for slide {i+1} of {len(data['slides'])}...</strong><br>
                            <small>Processing: "{slide['title']}" with {voice_style} voice style</small>
                        </div>
                        """, unsafe_allow_html=True)
                        
                        try:
                            speaker_notes = slide.get('speaker_notes', f"This is slide {i+1}: {slide.get('title', 'Untitled')}")
                            
                            # Add smart pauses if enabled
                            if st.session_state.advanced_features.get('include_pauses', True):
                                speaker_notes = speaker_notes.replace('. ', '. ... ')
                                speaker_notes = speaker_notes.replace('! ', '! ... ')
                                speaker_notes = speaker_notes.replace('? ', '? ... ')
                            
                            audio_content = lesson_gen.generate_enhanced_audio(
                                speaker_notes,
                                stability=settings['stability'],
                                similarity=settings['similarity']
                            )
                            
                            if audio_content:
                                audio_files.append((f"slide_{i+1:02d}_{slide.get('title', 'untitled').replace(' ', '_')[:20]}.mp3", audio_content))
                        except Exception as e:
                            st.warning(f"Error generating audio for slide {i+1}: {str(e)}")
                            continue
                    
                    # Enhanced completion status
                    status_container.markdown(f"""
                    <div class="status-success">
                        ✅ <strong>Enhanced generation complete!</strong><br>
                        Professional lesson materials with {theme_name} theme and {voice_style} narration are ready.
                        <br><small>Generated {len(data['slides'])} slides with {len(audio_files)} audio files</small>
                    </div>
                    """, unsafe_allow_html=True)
                    
                    # Store results
                    st.session_state.pptx_buffer = pptx_buffer
                    st.session_state.audio_files = audio_files
                    st.session_state.video_path = None
                    
                    # Add to generation history
                    generation_record = {
                        'timestamp': datetime.now().isoformat(),
                        'title': data['title'],
                        'theme': st.session_state.selected_theme,
                        'slide_count': len(data['slides']),
                        'audio_count': len(audio_files),
                        'success': True,
                        'research_enabled': bool(data.get('research_data')),
                        'voice_style': voice_style
                    }
                    st.session_state.generation_history.append(generation_record)
                    
                    st.session_state.current_step = 6
                    
                    time.sleep(2)
                    st.rerun()
                else:
                    st.markdown("""
                    <div style="background: linear-gradient(135deg, #fef2f2, #fecaca); padding: 2rem; border-radius: 16px; border-left: 4px solid #ef4444;">
                        ❌ <strong>PowerPoint generation failed.</strong><br>
                        Please try again or go back to regenerate the slides.
                    </div>
                    """, unsafe_allow_html=True)
                
                st.markdown('</div>', unsafe_allow_html=True)
            
            # Step 6: Enhanced Final Output
            elif st.session_state.current_step == 6:
                st.markdown(f"""
                <div class="status-success" style="text-align: center; padding: 3rem;">
                    <h2 style="color: #065f46; font-family: 'Space Grotesk', serif; margin-bottom: 1rem;">🎉 Your Enhanced Lesson is Ready!</h2>
                    <p style="font-size: 1.2rem;">Sophisticated, research-driven educational content with AI intelligence</p>
                </div>
                """, unsafe_allow_html=True)
                
                data = st.session_state.lesson_data
                
                # Enhanced summary metrics
                st.markdown("<h3 style='color: #8b5cf6; font-family: Inter, sans-serif; text-align: center; margin: 2rem 0;'>📊 Enhanced Lesson Summary</h3>", unsafe_allow_html=True)
                
                summary_cols = st.columns(6)
                
                metrics = [
                    ("Slides", len(data['slides']), "🖼️"),
                    ("Audio Files", len(st.session_state.audio_files) if hasattr(st.session_state, 'audio_files') else 0, "🎙️"),
                    ("Duration", f"{data['duration']} min", "⏱️"),
                    ("Research", "✅" if data.get('research_data') else "❌", "🔬"),
                    ("Theme", st.session_state.selected_theme.title(), "🎨"),
                    ("Voice Style", st.session_state.advanced_features.get('voice_style', 'Professional'), "🗣️")
                ]
                
                for i, (label, value, icon) in enumerate(metrics):
                    with summary_cols[i]:
                        st.markdown(f"""
                        <div class="metric-card">
                            <div class="metric-value">{icon}</div>
                            <div class="metric-label">{label}</div>
                            <small style="color: #6b7280;">{value}</small>
                        </div>
                        """, unsafe_allow_html=True)
                
                # Enhanced download section
                st.markdown("<h3 style='color: #14b8a6; font-family: Inter, sans-serif; text-align: center; margin: 2rem 0;'>📥 Download Enhanced Materials</h3>", unsafe_allow_html=True)
                
                download_cols = st.columns(4)
                
                with download_cols[0]:
                    if hasattr(st.session_state, 'pptx_buffer') and st.session_state.pptx_buffer:
                        filename = f"{data['title']}_{st.session_state.selected_theme}_enhanced.pptx"
                        st.download_button(
                            label="📄 Download PowerPoint",
                            data=st.session_state.pptx_buffer.getvalue(),
                            file_name=filename,
                            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                            help=f"Enhanced PowerPoint with {st.session_state.selected_theme} theme",
                            use_container_width=True
                        )
                
                with download_cols[1]:
                    if hasattr(st.session_state, 'audio_files') and st.session_state.audio_files:
                        # Create enhanced ZIP file
                        zip_buffer = io.BytesIO()
                        
                        with zipfile.ZipFile(zip_buffer, 'w', zipfile.ZIP_DEFLATED) as zip_file:
                            # Add audio files
                            for filename, audio_content in st.session_state.audio_files:
                                zip_file.writestr(f"audio/{filename}", audio_content)
                            
                            # Add lesson summary
                            summary_text = f"""
Enhanced Lesson Summary
======================
Title: {data['title']}
Subject: {data['subject']}
Grade Level: {data['grade_level']}
Duration: {data['duration']} minutes
Theme: {st.session_state.selected_theme}
Voice Style: {st.session_state.advanced_features.get('voice_style', 'Professional')}

Generated: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}
Slides: {len(data['slides'])}
Audio Files: {len(st.session_state.audio_files)}
Research Enabled: {'Yes' if data.get('research_data') else 'No'}

Learning Objectives:
{data['objectives']}
                            """
                            zip_file.writestr("lesson_summary.txt", summary_text)
                        
                        zip_buffer.seek(0)
                        
                        st.download_button(
                            label="🎙️ Download Audio Pack",
                            data=zip_buffer.getvalue(),
                            file_name=f"{data['title']}_enhanced_audio.zip",
                            mime="application/zip",
                            help="Enhanced audio narration with lesson summary",
                            use_container_width=True
                        )
                
                with download_cols[2]:
                    # Research notes download
                    if data.get('research_data'):
                        research_notes = f"""
Enhanced Research Notes
======================
Topic: {data['title']}
Research Depth: {data.get('research_depth', 'Standard')}
Generated: {data['research_data'].get('timestamp', 'Unknown')}

{data['research_data']['content']}

Enhanced Facts:
{data['facts']}

Lesson Outline:
{data['outline']}
                        """
                        
                        st.download_button(
                            label="🔬 Download Research",
                            data=research_notes,
                            file_name=f"{data['title']}_research_notes.txt",
                            mime="text/plain",
                            help="Comprehensive research and analysis",
                            use_container_width=True
                        )
                    else:
                        st.markdown("""
                        <div style="background: #f9fafb; padding: 1rem; border-radius: 12px; text-align: center; height: 60px; display: flex; align-items: center; justify-content: center;">
                            <span style="color: #64748b;">🔬 No Research Data</span>
                        </div>
                        """, unsafe_allow_html=True)
                
                with download_cols[3]:
                    st.markdown("""
                    <div style="background: linear-gradient(135deg, #f1f5f9, #e2e8f0); padding: 1rem; border-radius: 12px; text-align: center; height: 60px; display: flex; align-items: center; justify-content: center;">
                        <span style="color: #64748b; font-weight: 500;">🎬 Video: Coming Soon</span>
                    </div>
                    """, unsafe_allow_html=True)
                
                # Individual audio files
                if hasattr(st.session_state, 'audio_files') and st.session_state.audio_files:
                    with st.expander("🎵 Individual Audio Files", expanded=False):
                        st.markdown("<p style='text-align: center; color: #6b7280; margin-bottom: 1.5rem;'>Download individual slide narrations with enhanced voice styling:</p>", unsafe_allow_html=True)
                        
                        audio_cols = st.columns(3)
                        for i, (filename, audio_content) in enumerate(st.session_state.audio_files):
                            col_idx = i % 3
                            with audio_cols[col_idx]:
                                display_name = filename.replace('.mp3', '').replace('_', ' ').title()
                                st.download_button(
                                    label=f"🎙️ {display_name[:25]}...",
                                    data=audio_content,
                                    file_name=filename,
                                    mime="audio/mpeg",
                                    key=f"audio_{i}",
                                    help=f"Audio for {display_name}",
                                    use_container_width=True
                                )
                
                # Enhanced success message
                st.markdown(f"""
                <div class="status-success">
                    🎉 <strong>Enhanced lesson generation completed successfully!</strong><br>
                    Your sophisticated, research-driven materials with {st.session_state.selected_theme} theme are ready for the classroom.
                    <br><br>
                    <strong>Features Included:</strong>
                    • {len(data['slides'])} professionally designed slides
                    • {len(st.session_state.audio_files) if hasattr(st.session_state, 'audio_files') else 0} enhanced audio narrations
                    • {'Comprehensive research integration' if data.get('research_data') else 'Content-based generation'}
                    • Advanced {st.session_state.selected_theme} theme styling
                    • {st.session_state.advanced_features.get('voice_style', 'Professional')} voice narration
                </div>
                """, unsafe_allow_html=True)
                
                # Action buttons
                action_cols = st.columns(3)
                
                with action_cols[0]:
                    if st.button("🔄 Create New Lesson", type="primary", use_container_width=True):
                        # Reset session state for new lesson
                        keys_to_keep = ['generation_history', 'advanced_features']
                        for key in list(st.session_state.keys()):
                            if key not in keys_to_keep:
                                del st.session_state[key]
                        st.session_state.current_step = 1
                        st.rerun()
                
                with action_cols[1]:
                    if st.button("📊 View Analytics", type="secondary", use_container_width=True):
                        st.session_state.show_analytics = True
                        st.rerun()
                
                with action_cols[2]:
                    if st.button("💌 Share Feedback", type="secondary", use_container_width=True):
                        st.markdown("""
                        <div class="status-info">
                            <strong>💌 Love the enhanced features?</strong><br>
                            Your feedback helps us improve the AI-powered lesson generation experience!
                            <br><br>
                            <strong>What's new in this version:</strong><br>
                            • Advanced research integration<br>
                            • Enhanced theme options<br>
                            • Smart content suggestions<br>
                            • Professional voice styling<br>
                            • Comprehensive analytics
                        </div>
                        """, unsafe_allow_html=True)
    
    with tab2:
        render_advanced_features()
    
    with tab3:
        render_generation_analytics()

if __name__ == "__main__":
    main()
